#!/usr/bin/env -S uv run --script
# /// script
# requires-python = ">=3.11"
# dependencies = [
#   "python-pptx>=0.6.21",
#   "markitdown>=0.0.1a2",
#   "Pillow>=10.0.0",
# ]
# ///
"""
Parse PowerPoint files using 4 extraction methods: basic text+images, fast text-only,
detailed with notes, and comprehensive image extraction.

Usage: ./parse_pptx.py <file_path> <output_dir>
"""

import sys
import json
from pathlib import Path
from datetime import datetime
from pptx import Presentation
from markitdown import MarkItDown
from PIL import Image
import io


def parse_method1_python_pptx(source_file, output_dir):
    """Extract text and images using python-pptx."""
    print("    Method 1: python-pptx (basic)...")
    method_dir = output_dir / "python_pptx_basic"
    method_dir.mkdir(exist_ok=True)
    images_dir = method_dir / "images"
    images_dir.mkdir(exist_ok=True)

    try:
        prs = Presentation(source_file)

        text_content = []
        for slide_idx, slide in enumerate(prs.slides, 1):
            text_content.append(f"\n## Slide {slide_idx}\n")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text)

        with open(method_dir / "content.md", "w") as f:
            f.write("\n\n".join(text_content))

        img_count = 0
        for slide_idx, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    img_count += 1
                    ext = shape.image.content_type.split("/")[-1]
                    if ext == "jpeg":
                        ext = "jpg"
                    with open(images_dir / f"image_{img_count:03d}.{ext}", "wb") as f:
                        f.write(shape.image.blob)

        print(
            f"      ✓ Extracted text from {len(prs.slides)} slides and {img_count} images"
        )
        return True
    except Exception as e:
        print(f"      ⚠ Failed: {str(e)[:100]}")
        return False


def parse_method2_markitdown(source_file, output_dir):
    """Fast text-only extraction using Microsoft's markitdown."""
    print("    Method 2: markitdown...")
    method_dir = output_dir / "markitdown"
    method_dir.mkdir(exist_ok=True)

    try:
        md = MarkItDown()
        result = md.convert(str(source_file))

        with open(method_dir / "content.md", "w") as f:
            f.write(result.text_content)

        print(f"      ✓ Text extracted")
        return True
    except Exception as e:
        print(f"      ⚠ Failed: {str(e)[:100]}")
        return False


def parse_method3_python_pptx_detailed(source_file, output_dir):
    """Extract with speaker notes, layout info, and shape metadata."""
    print("    Method 3: python-pptx (detailed)...")
    method_dir = output_dir / "python_pptx_detailed"
    method_dir.mkdir(exist_ok=True)

    try:
        prs = Presentation(source_file)

        slides_data = []
        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_info = {
                "slide_number": slide_idx,
                "layout": slide.slide_layout.name
                if hasattr(slide.slide_layout, "name")
                else "Unknown",
                "shapes": [],
                "notes": "",
            }

            for shape in slide.shapes:
                shape_info = {
                    "type": shape.shape_type.name
                    if hasattr(shape.shape_type, "name")
                    else str(shape.shape_type),
                    "has_text": hasattr(shape, "text"),
                }
                if hasattr(shape, "text") and shape.text.strip():
                    shape_info["text"] = shape.text
                slide_info["shapes"].append(shape_info)

            if hasattr(slide, "notes_slide") and slide.notes_slide:
                try:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    if notes_text.strip():
                        slide_info["notes"] = notes_text
                except:
                    pass

            slides_data.append(slide_info)

        with open(method_dir / "slides_metadata.json", "w") as f:
            json.dump(slides_data, f, indent=2)

        md_content = []
        for slide_info in slides_data:
            md_content.append(
                f"\n## Slide {slide_info['slide_number']} ({slide_info['layout']})\n"
            )
            for shape in slide_info["shapes"]:
                if "text" in shape:
                    md_content.append(shape["text"])
            if slide_info["notes"]:
                md_content.append(f"\n**Notes:** {slide_info['notes']}\n")

        with open(method_dir / "content.md", "w") as f:
            f.write("\n\n".join(md_content))

        print(f"      ✓ Extracted detailed metadata for {len(slides_data)} slides")
        return True
    except Exception as e:
        print(f"      ⚠ Failed: {str(e)[:100]}")
        return False


def parse_method4_python_pptx_images(source_file, output_dir):
    """Extract all images with dimensions and metadata tracking."""
    print("    Method 4: python-pptx (images + metadata)...")
    method_dir = output_dir / "python_pptx_images"
    method_dir.mkdir(exist_ok=True)
    images_dir = method_dir / "images"
    images_dir.mkdir(exist_ok=True)

    try:
        prs = Presentation(source_file)

        images_metadata = []
        img_count = 0

        for slide_idx, slide in enumerate(prs.slides, 1):
            for shape_idx, shape in enumerate(slide.shapes):
                if hasattr(shape, "image"):
                    img_count += 1
                    ext = shape.image.content_type.split("/")[-1]
                    if ext == "jpeg":
                        ext = "jpg"

                    image_filename = (
                        f"slide_{slide_idx:02d}_image_{shape_idx:02d}.{ext}"
                    )
                    image_path = images_dir / image_filename

                    with open(image_path, "wb") as f:
                        f.write(shape.image.blob)

                    try:
                        img = Image.open(io.BytesIO(shape.image.blob))
                        width, height = img.size
                    except:
                        width, height = None, None

                    images_metadata.append(
                        {
                            "filename": image_filename,
                            "slide": slide_idx,
                            "shape_index": shape_idx,
                            "content_type": shape.image.content_type,
                            "width": width,
                            "height": height,
                            "size_bytes": len(shape.image.blob),
                        }
                    )

        with open(method_dir / "images_metadata.json", "w") as f:
            json.dump(images_metadata, f, indent=2)

        print(f"      ✓ Extracted {img_count} images with metadata")
        return True
    except Exception as e:
        print(f"      ⚠ Failed: {str(e)[:100]}")
        return False


def main():
    if len(sys.argv) != 3:
        print("Usage: ./parse_pptx.py <file_path> <output_dir>")
        sys.exit(1)

    source_file = Path(sys.argv[1])
    output_dir = Path(sys.argv[2])

    if not source_file.exists():
        print(f"Error: File not found: {source_file}")
        sys.exit(1)

    print(f"Parsing PPTX: {source_file.name}")
    print(f"Output directory: {output_dir}")
    print()

    output_dir.mkdir(parents=True, exist_ok=True)

    methods_success = {
        "method1_python_pptx_basic": parse_method1_python_pptx(source_file, output_dir),
        "method2_markitdown": parse_method2_markitdown(source_file, output_dir),
        "method3_python_pptx_detailed": parse_method3_python_pptx_detailed(
            source_file, output_dir
        ),
        "method4_python_pptx_images": parse_method4_python_pptx_images(
            source_file, output_dir
        ),
    }

    success_count = sum(methods_success.values())
    print(f"\n  ✓ Completed: {success_count}/{len(methods_success)} methods successful")

    results = {
        "timestamp": datetime.now().isoformat(),
        "source_file": str(source_file),
        "output_dir": str(output_dir),
        "methods": methods_success,
        "success_count": success_count,
    }

    with open(output_dir / "parsing_results.json", "w") as f:
        json.dump(results, f, indent=2)

    return 0 if success_count > 0 else 1


if __name__ == "__main__":
    sys.exit(main())
